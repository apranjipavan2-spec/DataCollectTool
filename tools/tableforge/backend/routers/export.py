import io
import json
import os
from datetime import datetime
from pathlib import Path
from typing import Optional

import numpy as np
import pandas as pd
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel

from ..shared import datasets, sanitize_for_json, EXPORTS_DIR, add_audit_log

router = APIRouter()


# ═══════════════════════════════════════════════════
# Module E: Export Engine
# ═══════════════════════════════════════════════════

class ExportConfig(BaseModel):
    dataset_id: str
    tables: list[dict]  # [{"name": "...", "headers": [...], "rows": [...], "title": "...", "subtitle": "..."}]
    format: str = "xlsx"
    filename: str = "TableForge_Export"
    options: dict = {}  # cover_page, cover_title, header_text, footer_text, page_numbers, include_raw_data, landscape


@router.post("/api/export")
async def export_tables(config: ExportConfig):
    dataset_required = config.format == "python" or config.options.get("include_raw_data") or config.options.get("formula_export")
    if dataset_required and config.dataset_id not in datasets:
        raise HTTPException(404, "Dataset not found — reload data to use raw data/formula export")

    if config.format == "xlsx":
        return await export_excel(config)
    elif config.format == "docx":
        return await export_word(config)
    elif config.format == "csv":
        return await export_csv(config)
    elif config.format == "pdf":
        return await export_pdf(config)
    elif config.format == "pptx":
        return await export_pptx(config)
    elif config.format == "python":
        return await export_python_script(config)
    else:
        raise HTTPException(400, f"Unsupported format: {config.format}")


async def export_pptx(config: ExportConfig):
    """One slide per table. Requires python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise HTTPException(500, "python-pptx not installed. pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    opts = config.options or {}
    cover_title = opts.get("cover_title") or config.filename

    # Cover slide
    if opts.get("cover_page", True):
        cover = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        title = cover.shapes.title
        title.text = cover_title
        if opts.get("cover_subtitle"):
            tb = cover.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
            tf = tb.text_frame
            tf.text = opts["cover_subtitle"]
            for p in tf.paragraphs:
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    for tbl in config.tables:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Title
        title_text = tbl.get("title") or tbl.get("name") or "Table"
        slide.shapes.title.text = title_text
        if tbl.get("subtitle"):
            sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.4))
            sub.text_frame.text = tbl["subtitle"]
            for p in sub.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        headers = tbl.get("headers") or []
        rows = tbl.get("rows") or []
        n_cols = len(headers)
        n_rows = len(rows) + 1
        if n_cols == 0 or n_rows == 1:
            continue

        # Cap to fit on slide
        max_rows = 18
        truncated = False
        if len(rows) > max_rows:
            rows = rows[:max_rows]
            n_rows = max_rows + 1
            truncated = True

        left, top = Inches(0.5), Inches(1.6)
        width, height = Inches(12.3), Inches(0.4 * n_rows)
        shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = shape.table

        # Header row
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(h)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x4A, 0x7B)

        # Data rows
        for i, row in enumerate(rows, start=1):
            for j in range(n_cols):
                cell = table.cell(i, j)
                val = row[j] if j < len(row) else ""
                cell.text = str(val) if val is not None else ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)

        if truncated:
            note = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
            note.text_frame.text = f"(showing first {max_rows} rows of {len(tbl.get('rows') or [])})"
            for p in note.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.italic = True
                p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    fname = f"{config.filename}.pptx"
    output_path = EXPORTS_DIR / fname
    prs.save(str(output_path))
    add_audit_log(config.dataset_id, "export_pptx", f"Exported {len(config.tables)} tables to PowerPoint")
    return {"path": str(output_path), "message": f"PowerPoint exported to {fname}", "download_filename": fname}


def _is_total_row(row):
    if not row:
        return False, False
    first = str(row[0])
    is_grand = first == "Grand Total"
    is_sub = not is_grand and ("Subtotal" in first or (len(row) > 1 and str(row[1]) == "Subtotal"))
    return is_grand, is_sub


def _build_number_format(vfmt):
    decimals = vfmt.get("decimals", 2)
    style = vfmt.get("style", "number")
    separator = vfmt.get("separator", True)
    prefix = vfmt.get("prefix", "")
    suffix = vfmt.get("suffix", "")
    currency = vfmt.get("currency_symbol", "$")

    int_part = "#,##0" if separator else "0"
    dec_part = "." + "0" * decimals if decimals > 0 else ""
    nf = int_part + dec_part

    if style == "percent":
        nf = nf + '"%"'
    elif style == "currency":
        nf = f'"{currency}"{nf}'
    if prefix:
        nf = f'"{prefix}"{nf}'
    if suffix:
        nf = f'{nf}"{suffix}"'
    return nf


async def export_excel(config: ExportConfig):
    """Export with full formatting using openpyxl."""
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    output_path = EXPORTS_DIR / f"{config.filename}.xlsx"

    from openpyxl import Workbook
    wb = Workbook()
    wb.remove(wb.active)

    thin_border = Border(
        left=Side(style="thin", color="D0D0D0"),
        right=Side(style="thin", color="D0D0D0"),
        top=Side(style="thin", color="D0D0D0"),
        bottom=Side(style="thin", color="D0D0D0"),
    )

    for t_idx, t in enumerate(config.tables):
        name = t.get("name", f"Table {t_idx + 1}")[:31]
        ws = wb.create_sheet(title=name)
        headers = t.get("headers", [])
        rows = t.get("rows", [])
        title = t.get("title", "")
        subtitle = t.get("subtitle", "")
        num_row_fields = t.get("num_row_fields", 1)
        cell_align = t.get("cell_align", "right")
        header_align = t.get("header_align", "center")
        row_label_align = t.get("row_label_align", "left")
        value_formats = t.get("value_formats", [])
        header_bg = (t.get("header_bg_color") or "1e293b").lstrip("#")
        header_tc = (t.get("header_text_color") or "FFFFFF").lstrip("#")
        total_bg = (t.get("total_bg_color") or "e8f0fe").lstrip("#")
        title_color = (t.get("title_color") or "")

        header_fill = PatternFill(start_color=header_bg or "1e293b", end_color=header_bg or "1e293b", fill_type="solid")
        header_font = Font(name="Segoe UI", size=11, bold=True, color=header_tc or "FFFFFF")
        data_font = Font(name="Segoe UI", size=10)
        total_fill = PatternFill(start_color=total_bg or "e8f0fe", end_color=total_bg or "e8f0fe", fill_type="solid")
        total_font = Font(name="Segoe UI", size=10, bold=True)
        subtotal_fill = PatternFill(start_color="f0f4f8", end_color="f0f4f8", fill_type="solid")

        # Per-column header format overrides
        hdr_formats_list = t.get("header_formats") or []
        hdr_fmt_map = {hf.get("field", ""): hf for hf in hdr_formats_list}

        def _get_header_font_fill(col_name):
            """Return (font, fill) for a header cell, applying per-column overrides."""
            fmt = hdr_fmt_map.get(col_name)
            if not fmt:
                return header_font, header_fill
            f_name = fmt.get("font") or "Segoe UI"
            f_size = fmt.get("size") or 11
            f_bold = fmt.get("bold") if fmt.get("bold") is not None else True
            f_italic = fmt.get("italic") or False
            f_color = (fmt.get("color") or header_tc or "FFFFFF").lstrip("#")
            f_bg = (fmt.get("backgroundColor") or header_bg or "1e293b").lstrip("#")
            font = Font(name=f_name, size=f_size, bold=f_bold, italic=f_italic, color=f_color)
            fill = PatternFill(start_color=f_bg, end_color=f_bg, fill_type="solid")
            return font, fill

        row_offset = 1
        if title:
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 1))
            cell = ws.cell(row=1, column=1, value=title)
            tc = title_color.lstrip("#") if title_color else ""
            cell.font = Font(name="Segoe UI", size=14, bold=True, color=tc if tc else "000000")
            row_offset = 2
        if subtitle:
            ws.merge_cells(start_row=row_offset, start_column=1, end_row=row_offset, end_column=max(len(headers), 1))
            cell = ws.cell(row=row_offset, column=1, value=subtitle)
            cell.font = Font(name="Segoe UI", size=11, color="666666")
            row_offset += 1
        if title or subtitle:
            row_offset += 1

        column_groups = t.get("column_groups")
        has_multi_level = column_groups and column_groups.get("has_multi_level")

        if has_multi_level:
            top_groups = column_groups["top"]
            bottom_labels = column_groups.get("bottom", [])
            # Row-dimension columns: merge across 2 rows
            for ci in range(1, num_row_fields + 1):
                ws.merge_cells(start_row=row_offset, start_column=ci, end_row=row_offset + 1, end_column=ci)
                col_name = headers[ci - 1]
                hf, hfl = _get_header_font_fill(col_name)
                cell = ws.cell(row=row_offset, column=ci, value=col_name)
                cell.font = hf
                cell.fill = hfl
                cell.alignment = Alignment(horizontal=header_align, vertical="center", wrap_text=True)
                cell.border = thin_border
            # Top-level group headers
            col_cursor = num_row_fields + 1
            for g in top_groups:
                colspan = g.get("colspan", 1)
                if colspan > 1:
                    ws.merge_cells(start_row=row_offset, start_column=col_cursor, end_row=row_offset, end_column=col_cursor + colspan - 1)
                cell = ws.cell(row=row_offset, column=col_cursor, value=g.get("label", ""))
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                cell.border = thin_border
                col_cursor += colspan
            # Bottom-level sub-headers
            for bi, label in enumerate(bottom_labels):
                ci = num_row_fields + 1 + bi
                col_name = headers[num_row_fields + bi] if (num_row_fields + bi) < len(headers) else label
                hf, hfl = _get_header_font_fill(col_name)
                cell = ws.cell(row=row_offset + 1, column=ci, value=label)
                cell.font = hf
                cell.fill = hfl
                cell.alignment = Alignment(horizontal=header_align, vertical="center", wrap_text=True)
                cell.border = thin_border
            row_offset += 1  # extra header row
        else:
            for ci, h in enumerate(headers, 1):
                hf, hfl = _get_header_font_fill(h)
                cell = ws.cell(row=row_offset, column=ci, value=h)
                cell.font = hf
                cell.fill = hfl
                cell.alignment = Alignment(horizontal=header_align, vertical="center", wrap_text=True)
                cell.border = thin_border

        ann_lookup = {}
        for ann in t.get("annotations", []):
            ann_lookup[(ann["rowIdx"], ann["colIdx"])] = ann.get("text", "")

        for ri, row in enumerate(rows, row_offset + 1):
            is_grand, is_sub = _is_total_row(row)
            data_ri = ri - row_offset - 1
            for ci, val in enumerate(row, 1):
                cell = ws.cell(row=ri, column=ci, value=val)
                cell.border = thin_border

                if is_grand:
                    cell.font = total_font
                    cell.fill = total_fill
                elif is_sub:
                    cell.font = total_font
                    cell.fill = subtotal_fill
                else:
                    cell.font = data_font

                is_val_col = (ci - 1) >= num_row_fields
                if isinstance(val, (int, float)):
                    align = cell_align if is_val_col else row_label_align
                    cell.alignment = Alignment(horizontal=align)
                    if is_val_col and value_formats:
                        vf_idx = (ci - 1 - num_row_fields) % len(value_formats)
                        cell.number_format = _build_number_format(value_formats[vf_idx])
                    else:
                        cell.number_format = '#,##0.00' if isinstance(val, float) else '#,##0'
                elif is_val_col:
                    cell.alignment = Alignment(horizontal=cell_align)
                else:
                    cell.alignment = Alignment(horizontal=row_label_align)

                ann_text = ann_lookup.get((data_ri, ci - 1))
                if ann_text:
                    from openpyxl.comments import Comment
                    comment = Comment(ann_text, "TableForge")
                    comment.width = 200
                    comment.height = 100
                    cell.comment = comment

        for ci in range(1, len(headers) + 1):
            max_len = max(
                (len(str(ws.cell(row=r, column=ci).value or "")) for r in range(row_offset, row_offset + len(rows) + 1)),
                default=10
            )
            ws.column_dimensions[get_column_letter(ci)].width = min(max(max_len + 2, 10), 40)

    # Formula-based sheet (additional sheet with SUM/AVERAGE formulas referencing raw data)
    if config.options.get("formula_export") and config.dataset_id in datasets:
        raw_df = datasets[config.dataset_id]["df"]
        formula_ws = wb.create_sheet(title="Formulas")
        formula_ws.cell(row=1, column=1, value="Formula Sheet — references Raw Data sheet")
        formula_ws.cell(row=1, column=1).font = Font(bold=True, name="Segoe UI")
        # Add SUM formulas for numeric columns
        numeric_cols = [c for c in raw_df.columns if str(raw_df[c].dtype) in ("int64", "float64")]
        formula_ws.cell(row=2, column=1, value="Column").font = Font(bold=True)
        formula_ws.cell(row=2, column=2, value="SUM").font = Font(bold=True)
        formula_ws.cell(row=2, column=3, value="AVERAGE").font = Font(bold=True)
        formula_ws.cell(row=2, column=4, value="COUNT").font = Font(bold=True)
        for r, col in enumerate(numeric_cols, 3):
            if "Raw Data" in [ws.title for ws in wb.worksheets]:
                col_letter = get_column_letter(list(raw_df.columns).index(col) + 1)
                n_rows = len(raw_df)
                formula_ws.cell(row=r, column=1, value=col)
                formula_ws.cell(row=r, column=2, value=f"='Raw Data'!{col_letter}2:{col_letter}{n_rows+1}")
                formula_ws.cell(row=r, column=3, value=f"=AVERAGE('Raw Data'!{col_letter}2:{col_letter}{n_rows+1})")
                formula_ws.cell(row=r, column=4, value=f"=COUNT('Raw Data'!{col_letter}2:{col_letter}{n_rows+1})")
            else:
                formula_ws.cell(row=r, column=1, value=col)
                formula_ws.cell(row=r, column=2, value=float(raw_df[col].sum()))
                formula_ws.cell(row=r, column=3, value=float(raw_df[col].mean()))
                formula_ws.cell(row=r, column=4, value=int(raw_df[col].count()))

    # Raw data sheet
    if config.options.get("include_raw_data") and config.dataset_id in datasets:
        raw_df = datasets[config.dataset_id]["df"]
        raw_ws = wb.create_sheet(title="Raw Data")
        raw_hdr_font = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
        raw_hdr_fill = PatternFill(start_color="1e293b", end_color="1e293b", fill_type="solid")
        for ci, col in enumerate(raw_df.columns, 1):
            cell = raw_ws.cell(row=1, column=ci, value=col)
            cell.font = raw_hdr_font
            cell.fill = raw_hdr_fill
        for ri, row in enumerate(raw_df.head(10000).values.tolist(), 2):
            for ci, val in enumerate(row, 1):
                raw_ws.cell(row=ri, column=ci, value=val if not (isinstance(val, float) and (np.isnan(val) or np.isinf(val))) else None)

    wb.save(output_path)
    fname = output_path.name
    return {"path": str(output_path), "message": f"Exported to {fname}", "download_filename": fname}


async def export_word(config: ExportConfig):
    """Export tables to a Word document with formatting matching preview."""
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt, Cm, Emu, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.section import WD_ORIENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import OxmlElement, parse_xml

    output_path = EXPORTS_DIR / f"{config.filename}.docx"
    doc = DocxDocument()
    opts = config.options or {}

    def _hex_to_rgb(hex_str):
        h = (hex_str or "").lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        return None

    def _set_cell_shading(cell, hex_color):
        c = (hex_color or "").lstrip("#")
        if not c or len(c) != 6:
            return
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{c}" w:val="clear"/>')
        cell._tc.get_or_add_tcPr().append(shading)

    def _set_cell_borders(cell, **kwargs):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcBorders = OxmlElement('w:tcBorders')
        for edge, attrs in kwargs.items():
            el = OxmlElement(f'w:{edge}')
            el.set(qn('w:val'), attrs.get('val', 'single'))
            el.set(qn('w:sz'), str(attrs.get('sz', 4)))
            el.set(qn('w:color'), attrs.get('color', '000000'))
            el.set(qn('w:space'), '0')
            tcBorders.append(el)
        tcPr.append(tcBorders)

    def _set_row_height(row, height_pt):
        tr = row._tr
        trPr = tr.get_or_add_trPr()
        trHeight = OxmlElement('w:trHeight')
        trHeight.set(qn('w:val'), str(int(height_pt * 20)))
        trHeight.set(qn('w:hRule'), 'atLeast')
        trPr.append(trHeight)

    def _set_cell_margin(cell, top=0, bottom=0, left=40, right=40):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcMar = OxmlElement('w:tcMar')
        for edge, val in [('top', top), ('bottom', bottom), ('start', left), ('end', right)]:
            el = OxmlElement(f'w:{edge}')
            el.set(qn('w:w'), str(val))
            el.set(qn('w:type'), 'dxa')
            tcMar.append(el)
        tcPr.append(tcMar)

    PX_PER_CHAR = 7.0

    def _col_pixel_widths(t_data):
        """Return list of pixel-scale weights for each column.
        Headers wrap in the preview, so cap their contribution.
        Content length is the primary width driver."""
        headers = t_data.get("headers", [])
        rows = t_data.get("rows", [])
        formatted_rows = t_data.get("formatted_rows") or []
        use_fmt = len(formatted_rows) == len(rows)
        col_widths_cfg = t_data.get("column_widths") or {}
        num_row_fields = t_data.get("num_row_fields", 1)
        px_list = []
        for ci, h in enumerate(headers):
            cfg_w = col_widths_cfg.get(h)
            if cfg_w and cfg_w > 0:
                px_list.append(float(cfg_w))
            else:
                content_max = 0
                sample = formatted_rows[:100] if use_fmt else rows[:100]
                for row in sample:
                    if ci < len(row):
                        val = str(row[ci] or "")
                        line_max = max((len(line) for line in val.split('\n')), default=0)
                        content_max = max(content_max, line_max)
                hdr_len = min(len(str(h)), 15)
                effective_len = max(content_max, hdr_len)
                est_px = max(effective_len * PX_PER_CHAR + 16, 45)
                if ci < num_row_fields:
                    est_px *= 1.1
                px_list.append(est_px)
        return px_list

    def _estimate_table_width_cm(t_data):
        return sum(_col_pixel_widths(t_data)) * 0.0264

    def _setup_section(section, landscape=False):
        section.top_margin = Cm(1.27)
        section.bottom_margin = Cm(1.27)
        section.left_margin = Cm(1.27)
        section.right_margin = Cm(1.27)
        if landscape:
            section.orientation = WD_ORIENT.LANDSCAPE
            if section.page_width < section.page_height:
                section.page_width, section.page_height = section.page_height, section.page_width
        else:
            section.orientation = WD_ORIENT.PORTRAIT
            if section.page_width > section.page_height:
                section.page_width, section.page_height = section.page_height, section.page_width

    def _setup_header_footer(section):
        if opts.get("header_text"):
            header = section.header
            header.is_linked_to_previous = False
            hp = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
            hp.text = ""
            run = hp.add_run(opts["header_text"])
            run.font.size = Pt(8)
            run.font.name = BODY_FONT
            run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        if opts.get("footer_text") or opts.get("page_numbers"):
            footer = section.footer
            footer.is_linked_to_previous = False
            fp = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            txt = opts.get("footer_text", "")
            if opts.get("page_numbers"):
                txt = (txt + "  " if txt else "") + "Page "
            if txt:
                run = fp.add_run(txt)
                run.font.size = Pt(8)
                run.font.name = BODY_FONT
            if opts.get("page_numbers"):
                fldChar = OxmlElement('w:fldChar')
                fldChar.set(qn('w:fldCharType'), 'begin')
                instrText = OxmlElement('w:instrText')
                instrText.text = "PAGE"
                fldChar2 = OxmlElement('w:fldChar')
                fldChar2.set(qn('w:fldCharType'), 'end')
                run = fp.add_run()
                run.font.size = Pt(8)
                run._r.append(fldChar)
                run._r.append(instrText)
                run._r.append(fldChar2)

    BODY_FONT = "Times New Roman"
    BODY_SIZE = 12
    HEADER_SIZE = 12
    CELL_MARGIN = 40

    portrait_usable_cm = 21.0 - 2.54
    landscape_usable_cm = 29.7 - 2.54

    def _embed_chart_if_present(t_data, is_landscape):
        """If the table carries a chart_image (base64 PNG), embed it after the table.
        The chart_title is rendered as a separate, editable Word paragraph above
        the image so the user can tweak the caption directly in Word."""
        chart_b64 = t_data.get("chart_image")
        if not chart_b64:
            return
        chart_title = t_data.get("chart_title") or ""
        try:
            import base64 as _b64
            img_bytes = _b64.b64decode(chart_b64)
        except Exception as decode_err:
            warn = doc.add_paragraph()
            wr = warn.add_run(f"[Chart image could not be decoded: {decode_err}]")
            wr.font.size = Pt(8); wr.italic = True
            wr.font.color.rgb = RGBColor(0xAA, 0x33, 0x33)
            return
        if chart_title:
            cp = doc.add_paragraph()
            cp.paragraph_format.space_before = Pt(8)
            cp.paragraph_format.space_after = Pt(3)
            cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            cr = cp.add_run(chart_title)
            cr.bold = True
            cr.font.size = Pt(11)
            cr.font.name = BODY_FONT
            cr.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        max_cm = (landscape_usable_cm if is_landscape else portrait_usable_cm) * 0.9
        target_cm = min(max_cm, 16.0)
        try:
            pic_p = doc.add_paragraph()
            pic_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            pic_p.paragraph_format.space_after = Pt(6)
            run = pic_p.add_run()
            run.add_picture(io.BytesIO(img_bytes), width=Cm(target_cm))
        except Exception as embed_err:
            warn = doc.add_paragraph()
            wr = warn.add_run(f"[Chart could not be embedded: {embed_err}]")
            wr.font.size = Pt(8); wr.italic = True
            wr.font.color.rgb = RGBColor(0xAA, 0x33, 0x33)

    # Cover page — use first section for it
    if opts.get("cover_page"):
        _setup_section(doc.sections[0], landscape=False)
        cover_p = doc.add_paragraph()
        cover_p.add_run("\n\n\n\n")
        cover_h = doc.add_heading(opts.get("cover_title") or config.filename, level=1)
        cover_h.runs[0].font.size = Pt(28)
        cover_h.runs[0].font.name = BODY_FONT
        if opts.get("cover_subtitle"):
            cs = doc.add_paragraph(opts["cover_subtitle"])
            cs.runs[0].font.size = Pt(14)
            cs.runs[0].font.name = BODY_FONT
            cs.runs[0].font.color.rgb = RGBColor(80, 80, 80)
        doc.add_paragraph(f"\nGenerated: {datetime.now().strftime('%B %d, %Y')}")
        _setup_header_footer(doc.sections[0])
        first_table_idx = 0
    else:
        first_table_idx = 0

    for t_idx, t in enumerate(config.tables):
        # Each table decides orientation based on its own content width
        est_width = _estimate_table_width_cm(t)
        table_landscape = est_width > portrait_usable_cm

        if t_idx == 0 and not opts.get("cover_page"):
            section = doc.sections[0]
            _setup_section(section, table_landscape)
            _setup_header_footer(section)
        else:
            # Honour per-table "Page Break Before" before starting a new section.
            if t.get("page_break_before"):
                doc.add_page_break()
            doc.add_section()
            section = doc.sections[-1]
            _setup_section(section, table_landscape)
            _setup_header_footer(section)

        title = t.get("title", t.get("name", f"Table {t_idx + 1}"))
        subtitle = t.get("subtitle", "")
        headers = t.get("headers", [])
        rows = t.get("rows", [])
        formatted_rows = t.get("formatted_rows") or []
        num_row_fields = t.get("num_row_fields", 1)
        cell_align = t.get("cell_align", "right")
        header_align = t.get("header_align", "center")
        row_label_align = t.get("row_label_align", "left")
        column_widths_cfg = t.get("column_widths") or {}
        row_height_cfg = t.get("row_height") or 0
        zebra = t.get("zebra", False)
        zebra_color = (t.get("zebra_color") or "#f8f9fa").lstrip("#")
        header_bg = (t.get("header_bg_color") or "1e293b").lstrip("#")
        header_tc = (t.get("header_text_color") or "FFFFFF").lstrip("#")
        total_bg = (t.get("total_bg_color") or "e8f0fe").lstrip("#")
        row_bg = (t.get("row_bg_color") or "").lstrip("#")
        table_bg = (t.get("bg_color") or "").lstrip("#")
        title_color = (t.get("title_color") or "").lstrip("#")

        align_map = {"left": WD_ALIGN_PARAGRAPH.LEFT, "center": WD_ALIGN_PARAGRAPH.CENTER, "right": WD_ALIGN_PARAGRAPH.RIGHT}

        hdr_formats_list = t.get("header_formats") or []
        hdr_fmt_map = {hf.get("field", ""): hf for hf in hdr_formats_list}

        column_groups = t.get("column_groups")
        has_multi_level = column_groups and column_groups.get("has_multi_level")
        group_start_cols = set()
        if has_multi_level:
            offset = num_row_fields
            for g in column_groups["top"]:
                if offset > num_row_fields:
                    group_start_cols.add(offset)
                offset += g.get("colspan", 1)

        # Title — matching preview font and size
        if title:
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.space_before = Pt(0)
            pf.space_after = Pt(4)
            run = p.add_run(title)
            run.font.name = BODY_FONT
            title_size_px = t.get("title_size") or 18
            run.font.size = Pt(round(title_size_px * 0.75))
            run.bold = t.get("title_bold", True)
            run.italic = t.get("title_italic", False)
            if title_color:
                run.font.color.rgb = _hex_to_rgb(title_color) or RGBColor(0, 0, 0)
            p.alignment = align_map.get(t.get("title_align", "left"), WD_ALIGN_PARAGRAPH.LEFT)
        if subtitle:
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.space_before = Pt(0)
            pf.space_after = Pt(6)
            run = p.add_run(subtitle)
            run.font.size = Pt(BODY_SIZE)
            run.font.name = BODY_FONT
            run.font.color.rgb = RGBColor(100, 100, 100)

        if not headers:
            # Chart-only item: no table to render, but still embed the chart
            # (e.g. the user picked a chart in the export dialog without its
            # source table). Falls through silently if no chart_image is set.
            if t.get("chart_image"):
                _embed_chart_if_present(t, table_landscape)
            continue

        styled_html = t.get("styled_html", "")
        if styled_html:
            from htmldocx import HtmlToDocx
            parser = HtmlToDocx()
            tables_before = len(doc.tables)
            parser.add_html_to_document(styled_html, doc)
            if len(doc.tables) > tables_before:
                tbl = doc.tables[-1]
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                tbl.autofit = False
                tbl_el = tbl._tbl
                tblPr = tbl_el.tblPr if tbl_el.tblPr is not None else OxmlElement('w:tblPr')
                tblW = OxmlElement('w:tblW')
                tblW.set(qn('w:type'), 'auto')
                tblW.set(qn('w:w'), '0')
                for old in tblPr.findall(qn('w:tblW')):
                    tblPr.remove(old)
                tblPr.append(tblW)
                layout = OxmlElement('w:tblLayout')
                layout.set(qn('w:type'), 'fixed')
                for old in tblPr.findall(qn('w:tblLayout')):
                    tblPr.remove(old)
                tblPr.append(layout)
                page_width_emu = section.page_width - section.left_margin - section.right_margin
                col_px = _col_pixel_widths(t)
                total_px = sum(col_px) or 1
                col_emus = [int(page_width_emu * px / total_px) for px in col_px]
                for ci in range(min(len(tbl.columns), len(col_emus))):
                    tbl.columns[ci].width = col_emus[ci]
                    for row in tbl.rows:
                        row.cells[ci].width = col_emus[ci]

                if has_multi_level and len(tbl.rows) >= 2:
                    for ci in range(num_row_fields):
                        tbl.rows[0].cells[ci].merge(tbl.rows[1].cells[ci])
                    col_cursor = num_row_fields
                    for g in column_groups["top"]:
                        colspan = g.get("colspan", 1)
                        if colspan > 1:
                            tbl.rows[0].cells[col_cursor].merge(tbl.rows[0].cells[col_cursor + colspan - 1])
                        col_cursor += colspan

                header_row_count = 2 if has_multi_level else 1
                for ri, row in enumerate(tbl.rows):
                    is_header = ri < header_row_count
                    for ci, cell in enumerate(row.cells):
                        is_val_col = ci >= num_row_fields
                        if is_header:
                            h_align = align_map.get(header_align, WD_ALIGN_PARAGRAPH.CENTER)
                        elif is_val_col:
                            h_align = align_map.get(cell_align, WD_ALIGN_PARAGRAPH.RIGHT)
                        else:
                            h_align = align_map.get(row_label_align, WD_ALIGN_PARAGRAPH.LEFT)
                        for p in cell.paragraphs:
                            p.alignment = h_align
                        cell.vertical_alignment = 1  # WD_CELL_VERTICAL_ALIGNMENT.CENTER

            footnote_text = t.get("footnote", "")
            footnotes_list = t.get("footnotes") or []
            if footnote_text:
                fp = doc.add_paragraph()
                fp.paragraph_format.space_before = Pt(4)
                fr = fp.add_run(footnote_text)
                fr.font.size = Pt(8)
                fr.font.name = BODY_FONT
                fr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                fr.italic = True
            for fn in footnotes_list:
                fp = doc.add_paragraph()
                fr = fp.add_run(f"{fn.get('marker', '*')} {fn.get('text', '')}")
                fr.font.size = Pt(8)
                fr.font.name = BODY_FONT
                fr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

            ann_list = [ann.get("text", "") for ann in t.get("annotations", [])]
            if ann_list:
                fn_para = doc.add_paragraph()
                fn_run = fn_para.add_run("Annotations:")
                fn_run.bold = True
                fn_run.font.size = Pt(8)
                fn_run.font.name = BODY_FONT
                for i, ann_text in enumerate(ann_list, 1):
                    fp = doc.add_paragraph(f"[{i}] {ann_text}")
                    fp.runs[0].font.size = Pt(8)
                    fp.runs[0].font.name = BODY_FONT
                    fp.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

            interpretation = t.get("interpretation", "")
            if interpretation:
                doc.add_paragraph()
                ip_heading = doc.add_paragraph()
                ip_run = ip_heading.add_run("Interpretation")
                ip_run.bold = True
                ip_run.font.size = Pt(11)
                ip_run.font.name = BODY_FONT
                ip_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                ip = doc.add_paragraph(interpretation)
                for run in ip.runs:
                    run.font.size = Pt(10)
                    run.font.name = BODY_FONT
                    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

            _embed_chart_if_present(t, table_landscape)
            continue

        header_row_count = 2 if has_multi_level else 1

        table = doc.add_table(rows=header_row_count + len(formatted_rows or rows), cols=len(headers))
        table.style = "Table Grid"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.autofit = False

        # Column widths: use pixel-scale weights (same scale as preview)
        page_width_emu = section.page_width - section.left_margin - section.right_margin
        col_px = _col_pixel_widths(t)

        total_px = sum(col_px) or 1
        for ci in range(len(headers)):
            if ci < len(table.columns) and ci < len(col_px):
                table.columns[ci].width = int(page_width_emu * col_px[ci] / total_px)

        def _style_header_cell(cell, text, align=header_align, col_name=None):
            cell.text = ""
            p = cell.paragraphs[0]
            p.alignment = align_map.get(align, WD_ALIGN_PARAGRAPH.CENTER)
            pf = p.paragraph_format
            pf.space_before = Pt(1)
            pf.space_after = Pt(1)
            pf.line_spacing = Pt(HEADER_SIZE + 2)
            run = p.add_run(str(text))
            fmt = hdr_fmt_map.get(col_name) if col_name else None
            if fmt:
                run.bold = fmt.get("bold", True)
                run.italic = fmt.get("italic", False)
                fmt_size_px = fmt.get("size")
                run.font.size = Pt(round(fmt_size_px * 0.75)) if fmt_size_px else Pt(HEADER_SIZE)
                run.font.name = fmt.get("font") or BODY_FONT
                fc = (fmt.get("color") or header_tc).lstrip("#")
                run.font.color.rgb = _hex_to_rgb(fc) or RGBColor(0xFF, 0xFF, 0xFF)
                bg = (fmt.get("backgroundColor") or header_bg).lstrip("#")
                _set_cell_shading(cell, bg)
            else:
                run.bold = True
                run.font.size = Pt(HEADER_SIZE)
                run.font.name = BODY_FONT
                run.font.color.rgb = _hex_to_rgb(header_tc) or RGBColor(0xFF, 0xFF, 0xFF)
                _set_cell_shading(cell, header_bg)
            _set_cell_margin(cell, top=20, bottom=20, left=CELL_MARGIN, right=CELL_MARGIN)
            cell.vertical_alignment = 1  # CENTER

        if has_multi_level:
            top_groups = column_groups["top"]
            bottom_labels = column_groups.get("bottom", [])
            for ci in range(num_row_fields):
                table.rows[0].cells[ci].merge(table.rows[1].cells[ci])
                _style_header_cell(table.rows[0].cells[ci], headers[ci], col_name=headers[ci])

            col_cursor = num_row_fields
            for gi, g in enumerate(top_groups):
                colspan = g.get("colspan", 1)
                if colspan > 1:
                    table.rows[0].cells[col_cursor].merge(table.rows[0].cells[col_cursor + colspan - 1])
                cell = table.rows[0].cells[col_cursor]
                _style_header_cell(cell, g.get("label", ""), "center")
                if gi > 0:
                    _set_cell_borders(cell, left={'sz': 12, 'val': 'single', 'color': '475569'})
                col_cursor += colspan

            for bi, label in enumerate(bottom_labels):
                ci = num_row_fields + bi
                if ci < len(headers):
                    col_name = headers[ci] if ci < len(headers) else label
                    _style_header_cell(table.rows[1].cells[ci], label, col_name=col_name)
                    if ci in group_start_cols:
                        _set_cell_borders(table.rows[1].cells[ci], left={'sz': 12, 'val': 'single', 'color': '475569'})
        else:
            for ci, h in enumerate(headers):
                _style_header_cell(table.rows[0].cells[ci], h, col_name=h)

        for hr_idx in range(header_row_count):
            _set_row_height(table.rows[hr_idx], HEADER_SIZE + 6)

        ann_lookup = {}
        ann_list = []
        for ann in t.get("annotations", []):
            key = (ann["rowIdx"], ann["colIdx"])
            ann_lookup[key] = len(ann_list) + 1
            ann_list.append(ann.get("text", ""))

        use_formatted = len(formatted_rows) == len(rows)
        for ri in range(len(rows)):
            raw_row = rows[ri]
            fmt_row = formatted_rows[ri] if use_formatted else None
            is_grand, is_sub = _is_total_row(raw_row)
            tbl_row = table.rows[ri + header_row_count]

            if row_height_cfg and row_height_cfg > 0:
                _set_row_height(tbl_row, row_height_cfg * 0.75)

            for ci in range(len(headers)):
                cell = tbl_row.cells[ci]
                if fmt_row and ci < len(fmt_row):
                    display = fmt_row[ci]
                elif ci < len(raw_row):
                    val = raw_row[ci]
                    if val is None:
                        display = ""
                    elif isinstance(val, float):
                        display = f"{val:,.2f}"
                    elif isinstance(val, int):
                        display = f"{val:,}"
                    else:
                        display = str(val)
                else:
                    display = ""

                ann_num = ann_lookup.get((ri, ci))
                if ann_num:
                    display = f"{display} [{ann_num}]"

                cell.text = ""
                p = cell.paragraphs[0]
                pf = p.paragraph_format
                pf.space_before = Pt(1)
                pf.space_after = Pt(1)
                pf.line_spacing = Pt(BODY_SIZE + 2)
                is_val_col = ci >= num_row_fields
                p.alignment = align_map.get(cell_align if is_val_col else row_label_align, WD_ALIGN_PARAGRAPH.LEFT)

                run = p.add_run(str(display))
                run.font.size = Pt(BODY_SIZE)
                run.font.name = BODY_FONT
                _set_cell_margin(cell, top=20, bottom=20, left=CELL_MARGIN, right=CELL_MARGIN)

                if is_grand:
                    run.bold = True
                    _set_cell_shading(cell, total_bg)
                elif is_sub:
                    run.bold = True
                    _set_cell_shading(cell, "f0f4f8")
                elif zebra and ri % 2 == 1:
                    _set_cell_shading(cell, zebra_color)
                elif not is_val_col and row_bg:
                    _set_cell_shading(cell, row_bg)
                elif is_val_col and table_bg:
                    _set_cell_shading(cell, table_bg)

                if ci in group_start_cols:
                    _set_cell_borders(cell, left={'sz': 12, 'val': 'single', 'color': '475569'})

        # Footnotes
        footnote_text = t.get("footnote", "")
        footnotes_list = t.get("footnotes") or []
        if footnote_text:
            fp = doc.add_paragraph()
            fp.paragraph_format.space_before = Pt(4)
            fr = fp.add_run(footnote_text)
            fr.font.size = Pt(8)
            fr.font.name = BODY_FONT
            fr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            fr.italic = True
        for fn in footnotes_list:
            fp = doc.add_paragraph()
            fr = fp.add_run(f"{fn.get('marker', '*')} {fn.get('text', '')}")
            fr.font.size = Pt(8)
            fr.font.name = BODY_FONT
            fr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        if ann_list:
            fn_para = doc.add_paragraph()
            fn_run = fn_para.add_run("Annotations:")
            fn_run.bold = True
            fn_run.font.size = Pt(8)
            fn_run.font.name = BODY_FONT
            for i, ann_text in enumerate(ann_list, 1):
                fp = doc.add_paragraph(f"[{i}] {ann_text}")
                fp.runs[0].font.size = Pt(8)
                fp.runs[0].font.name = BODY_FONT
                fp.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        interpretation = t.get("interpretation", "")
        if interpretation:
            doc.add_paragraph()
            ip_heading = doc.add_paragraph()
            ip_run = ip_heading.add_run("Interpretation")
            ip_run.bold = True
            ip_run.font.size = Pt(11)
            ip_run.font.name = BODY_FONT
            ip_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            ip = doc.add_paragraph(interpretation)
            for run in ip.runs:
                run.font.size = Pt(10)
                run.font.name = BODY_FONT
                run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        _embed_chart_if_present(t, table_landscape)

    doc.save(output_path)
    fname = output_path.name
    return {"path": str(output_path), "message": f"Exported to {fname}", "download_filename": fname}


async def export_csv(config: ExportConfig):
    if config.tables:
        t = config.tables[0]
        output_path = EXPORTS_DIR / f"{config.filename}.csv"
        headers = t.get("headers", [])
        rows = t.get("rows", [])
        column_groups = t.get("column_groups")
        has_multi_level = column_groups and column_groups.get("has_multi_level")
        num_row_fields = t.get("num_row_fields", 1)

        with open(output_path, "w", encoding="utf-8", newline="") as f:
            import csv
            writer = csv.writer(f)
            if has_multi_level:
                top_groups = column_groups["top"]
                bottom_labels = column_groups.get("bottom", [])
                top_row = list(headers[:num_row_fields])
                for g in top_groups:
                    top_row.append(g.get("label", ""))
                    top_row.extend([""] * (g.get("colspan", 1) - 1))
                writer.writerow(top_row)
                bot_row = [""] * num_row_fields + bottom_labels
                writer.writerow(bot_row)
            else:
                writer.writerow(headers)
            for row in rows:
                writer.writerow(row)

        fname = output_path.name
        return {"path": str(output_path), "message": f"Exported to {fname}", "download_filename": fname}
    raise HTTPException(400, "No tables to export")


async def export_pdf(config: ExportConfig):
    """Export tables as PDF using HTML+CSS rendering."""
    output_path = EXPORTS_DIR / f"{config.filename}.pdf"
    opts = config.options or {}

    # Build HTML
    html_parts = ["""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
    body { font-family: Segoe UI, Arial, sans-serif; font-size: 11px; margin: 20px; }
    h1 { font-size: 20px; margin-bottom: 4px; }
    h2 { font-size: 14px; margin-bottom: 2px; color: #333; }
    .subtitle { font-size: 11px; color: #666; margin-bottom: 12px; }
    .footnote { font-size: 10px; color: #666; margin-top: 4px; font-style: italic; }
    table { border-collapse: collapse; width: 100%; margin-bottom: 24px; }
    th { background: #1e293b; color: #fff; padding: 6px 10px; text-align: left; font-size: 10px; }
    td { padding: 5px 10px; border-bottom: 1px solid #ddd; font-size: 10px; }
    tr.grand-total td { background: #e8f0fe; font-weight: bold; }
    tr.subtotal td { background: #f0f4f8; font-weight: bold; }
    tr:nth-child(even) td { background: #f9f9f9; }
    .cover { text-align: center; padding-top: 80px; page-break-after: always; }
    @media print { .cover { page-break-after: always; } }
    </style></head><body>"""]

    if opts.get("cover_page"):
        cover_title = opts.get("cover_title") or config.filename
        cover_sub = opts.get("cover_subtitle", "")
        html_parts.append(f'<div class="cover"><h1>{cover_title}</h1>')
        if cover_sub: html_parts.append(f'<p class="subtitle">{cover_sub}</p>')
        html_parts.append(f'<p style="color:#999;margin-top:40px">{datetime.now().strftime("%B %d, %Y")}</p></div>')

    if opts.get("header_text"):
        html_parts.append(f'<div style="text-align:center;color:#666;margin-bottom:16px">{opts["header_text"]}</div>')

    for t in config.tables:
        title = t.get("title", t.get("name", ""))
        subtitle = t.get("subtitle", "")
        footnote = t.get("footnote", "")
        headers = t.get("headers", [])
        rows = t.get("rows", [])
        formatted_rows = t.get("formatted_rows") or []
        num_row_fields = t.get("num_row_fields", 1)
        cell_align = t.get("cell_align", "right")
        header_align = t.get("header_align", "center")
        row_label_align = t.get("row_label_align", "left")
        use_formatted = len(formatted_rows) == len(rows)

        if title: html_parts.append(f"<h2>{title}</h2>")
        if subtitle: html_parts.append(f'<p class="subtitle">{subtitle}</p>')
        if not headers: continue

        column_groups = t.get("column_groups")
        has_multi_level = column_groups and column_groups.get("has_multi_level")

        # Per-column header format overrides for PDF
        hdr_formats_list = t.get("header_formats") or []
        hdr_fmt_map = {hf.get("field", ""): hf for hf in hdr_formats_list}

        def _th_style(col_name, align=header_align):
            fmt = hdr_fmt_map.get(col_name)
            parts = [f"text-align:{align}"]
            if fmt:
                if fmt.get("font"): parts.append(f"font-family:{fmt['font']}")
                if fmt.get("size"): parts.append(f"font-size:{fmt['size']}px")
                if fmt.get("bold") is False: parts.append("font-weight:normal")
                if fmt.get("italic"): parts.append("font-style:italic")
                if fmt.get("color"): parts.append(f"color:{fmt['color']}")
                if fmt.get("backgroundColor"): parts.append(f"background:{fmt['backgroundColor']}")
            return ";".join(parts)

        html_parts.append("<table><thead>")
        if has_multi_level:
            top_groups = column_groups["top"]
            bottom_labels = column_groups.get("bottom", [])
            html_parts.append("<tr>")
            for ci in range(num_row_fields):
                style = _th_style(headers[ci], header_align)
                html_parts.append(f'<th rowspan="2" style="{style}">{headers[ci]}</th>')
            for g in top_groups:
                colspan = g.get("colspan", 1)
                html_parts.append(f'<th colspan="{colspan}" style="text-align:center">{g.get("label", "")}</th>')
            html_parts.append("</tr><tr>")
            for bi, label in enumerate(bottom_labels):
                col_name = headers[num_row_fields + bi] if (num_row_fields + bi) < len(headers) else label
                style = _th_style(col_name, header_align)
                html_parts.append(f'<th style="{style}">{label}</th>')
            html_parts.append("</tr>")
        else:
            html_parts.append("<tr>")
            for h in headers:
                style = _th_style(h, header_align)
                html_parts.append(f'<th style="{style}">{h}</th>')
            html_parts.append("</tr>")
        html_parts.append("</thead><tbody>")

        for ri, row in enumerate(rows):
            is_grand, is_sub = _is_total_row(row)
            cls = ' class="grand-total"' if is_grand else (' class="subtotal"' if is_sub else "")
            html_parts.append(f"<tr{cls}>")
            fmt_row = formatted_rows[ri] if use_formatted and ri < len(formatted_rows) else None
            for ci in range(len(headers)):
                if fmt_row and ci < len(fmt_row):
                    display = fmt_row[ci]
                elif ci < len(row):
                    cell = row[ci]
                    if cell is None:
                        display = ""
                    elif isinstance(cell, float):
                        display = f"{cell:,.2f}"
                    elif isinstance(cell, int):
                        display = f"{cell:,}"
                    else:
                        display = str(cell)
                else:
                    display = ""
                is_val_col = ci >= num_row_fields
                align = cell_align if is_val_col else row_label_align
                bold = "font-weight:bold;" if (is_grand or is_sub) else ""
                html_parts.append(f'<td style="text-align:{align};{bold}">{display}</td>')
            html_parts.append("</tr>")

        html_parts.append("</tbody></table>")
        if footnote: html_parts.append(f'<p class="footnote">{footnote}</p>')
        interpretation = t.get("interpretation", "")
        if interpretation:
            html_parts.append(f'<p style="margin-top:12px;"><strong>Interpretation</strong></p>')
            html_parts.append(f'<p style="font-family:Georgia,serif;font-size:11px;color:#444;line-height:1.6;">{interpretation}</p>')

    if opts.get("footer_text"):
        html_parts.append(f'<div style="text-align:center;color:#999;margin-top:20px;font-size:10px">{opts["footer_text"]}</div>')

    html_parts.append("</body></html>")
    html_content = "".join(html_parts)

    # Try weasyprint, fall back to HTML file
    try:
        import weasyprint
        weasyprint.HTML(string=html_content).write_pdf(str(output_path))
    except ImportError:
        # Save as HTML with .pdf extension note
        html_path = EXPORTS_DIR / f"{config.filename}.html"
        html_path.write_text(html_content, encoding="utf-8")
        output_path = html_path
        fname = html_path.name
        return {"path": str(output_path), "message": f"PDF (as HTML) saved to {fname} — open and print to PDF", "download_filename": fname}

    fname = output_path.name
    return {"path": str(output_path), "message": f"PDF exported to {fname}", "download_filename": fname}


async def export_python_script(config: ExportConfig):
    """Generate a pandas Python script that reproduces the tables."""
    ds = datasets[config.dataset_id]
    source_file = ds.get("filepath", "your_file.xlsx")
    filename_base = Path(source_file).name

    lines = [
        "#!/usr/bin/env python3",
        '"""',
        f"TableForge Export — Generated script",
        f"Source file: {filename_base}",
        '"""',
        "",
        "import pandas as pd",
        "import numpy as np",
        "",
        f"# Load source data",
        f"df = pd.read_excel({repr(filename_base)})  # adjust path as needed",
        f"# df = pd.read_csv({repr(filename_base)})  # for CSV files",
        "",
    ]

    for t_idx, t in enumerate(config.tables):
        name = t.get("name", f"Table {t_idx + 1}")
        headers = t.get("headers", [])
        rows = t.get("rows", [])
        lines.append(f"# {'─' * 50}")
        lines.append(f"# {name}")
        lines.append(f"# {'─' * 50}")
        lines.append(f"# Headers: {headers}")
        lines.append(f"table_{t_idx + 1}_data = [")
        for row in rows[:5]:
            lines.append(f"    {row},")
        if len(rows) > 5:
            lines.append(f"    # ... {len(rows) - 5} more rows")
        lines.append(f"]")
        lines.append(f"table_{t_idx + 1} = pd.DataFrame(table_{t_idx + 1}_data, columns={headers})")
        lines.append(f"print('\\n{name}:')")
        lines.append(f"print(table_{t_idx + 1}.to_string(index=False))")
        lines.append("")

    lines.append("# Save to Excel")
    lines.append(f"with pd.ExcelWriter('{config.filename}.xlsx', engine='openpyxl') as writer:")
    for t_idx, t in enumerate(config.tables):
        name = t.get("name", f"Table {t_idx + 1}")
        lines.append(f"    table_{t_idx + 1}.to_excel(writer, sheet_name={repr(name[:31])}, index=False)")
    lines.append("")
    lines.append("print('\\nSaved to {}.xlsx'.format(repr(config.filename)))")

    script_content = "\n".join(lines)
    output_path = EXPORTS_DIR / f"{config.filename}.py"
    output_path.write_text(script_content, encoding="utf-8")

    add_audit_log(config.dataset_id, "export", f"Python script exported: {config.filename}.py")
    return {"message": "Python script generated", "download_filename": f"{config.filename}.py"}


# ═══════════════════════════════════════════════════
# Download exported file
# ═══════════════════════════════════════════════════

@router.get("/api/export/download/{filename}")
async def download_export(filename: str):
    filepath = EXPORTS_DIR / filename
    if not filepath.exists():
        raise HTTPException(404, "File not found")
    media_types = {
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".csv": "text/csv",
        ".pdf": "application/pdf",
        ".py": "text/x-python",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    ext = filepath.suffix
    return FileResponse(filepath, media_type=media_types.get(ext, "application/octet-stream"), filename=filename)


# ═══════════════════════════════════════════════════
# Module J: Report Export
# ═══════════════════════════════════════════════════

class ReportExportConfig(BaseModel):
    dataset_id: str = ""
    elements: list[dict] = []
    filename: str = "TableForge_Report"


@router.post("/api/report/export")
async def export_report(config: ReportExportConfig):
    """Export a full structured report as Word document."""
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    output_path = EXPORTS_DIR / f"{config.filename}.docx"
    doc = DocxDocument()

    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    now_str = datetime.now().strftime("%B %d, %Y")

    for el in config.elements:
        etype = el.get("type", "text")
        text = str(el.get("text", "")).replace("{date}", now_str)

        if etype == "cover":
            doc.add_paragraph("\n\n")
            h = doc.add_heading(text or "Report", level=1)
            h.runs[0].font.size = Pt(28)
            doc.add_paragraph(f"\n{now_str}")
            doc.add_page_break()
        elif etype == "heading":
            level = el.get("level", 1)
            doc.add_heading(text, level=level)
        elif etype == "text":
            if text:
                doc.add_paragraph(text)
        elif etype == "page_break":
            doc.add_page_break()
        elif etype == "table":
            name = el.get("name", "")
            title = el.get("title", name)
            subtitle = el.get("subtitle", "")
            headers = el.get("headers", [])
            rows = el.get("rows", [])

            if title:
                tp = doc.add_heading(title, level=3)
            if subtitle:
                sp = doc.add_paragraph(subtitle)
                if sp.runs:
                    sp.runs[0].font.size = Pt(10)
                    sp.runs[0].font.color.rgb = RGBColor(80, 80, 80)

            if headers:
                tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                tbl.style = "Table Grid"
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

                for ci, h in enumerate(headers):
                    cell = tbl.rows[0].cells[ci]
                    cell.text = str(h)
                    p = cell.paragraphs[0]
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = p.runs[0] if p.runs else p.add_run(str(h))
                    run.bold = True
                    run.font.size = Pt(9)

                for ri, row in enumerate(rows):
                    is_total = len(row) > 0 and str(row[0]) == "Grand Total"
                    for ci, val in enumerate(row):
                        cell = tbl.rows[ri + 1].cells[ci]
                        display = f"{val:,.2f}" if isinstance(val, float) else f"{val:,}" if isinstance(val, int) else str(val) if val is not None else ""
                        cell.text = display
                        p = cell.paragraphs[0]
                        if isinstance(val, (int, float)):
                            p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        run = p.runs[0] if p.runs else p.add_run(display)
                        run.font.size = Pt(9)
                        if is_total:
                            run.bold = True

    doc.save(output_path)
    fname = output_path.name
    return {"path": str(output_path), "message": f"Report exported to {fname}", "download_filename": fname}
